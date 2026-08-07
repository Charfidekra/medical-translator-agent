import io
import gc
import os
import shutil
import numpy as np
import streamlit as st
import streamlit.components.v1 as components
from PIL import Image, ImageEnhance, ImageOps
from pptx import Presentation

# ----------------------------------------------------
# الاستيرادات مع معالجة الاستثناءات الآمنة
# ----------------------------------------------------
try:
    import fitz  # PyMuPDF
except ImportError:
    st.error("مكتبة PyMuPDF غير مثبتة في requirements.txt")

try:
    import cv2
except ImportError:
    cv2 = None

try:
    import easyocr
except ImportError:
    easyocr = None

try:
    import pytesseract
    tesseract_cmd_path = shutil.which("tesseract")
    if tesseract_cmd_path:
        pytesseract.pytesseract.tesseract_cmd = tesseract_cmd_path
except ImportError:
    pytesseract = None

from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
from main import translate_document


# ----------------------------------------------------
# 0. تهيئة محرك EasyOCR بأمان
# ----------------------------------------------------
@st.cache_resource(show_spinner=False)
def load_easyocr_reader():
    """تحميل EasyOCR عند الحاجة فقط دون إسقاط التطبيق"""
    if easyocr is not None:
        try:
            return easyocr.Reader(['fr', 'en'], gpu=False)
        except Exception:
            return None
    return None


# ----------------------------------------------------
# 1. واجهة الأذكار والتسبيحات أثناء المعالجة
# ----------------------------------------------------
AZKAR_HTML = """
<!DOCTYPE html>
<html lang="ar" dir="rtl">
<head>
<meta charset="utf-8">
<style>
  body {
    background-color: #0e1117;
    color: #ffffff;
    font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
    display: flex;
    flex-direction: column;
    align-items: center;
    justify-content: center;
    padding: 15px;
    margin: 0;
  }
  .card {
    background: #1e2530;
    border: 1px solid #313d4f;
    border-radius: 12px;
    padding: 20px;
    width: 90%;
    max-width: 450px;
    text-align: center;
    box-shadow: 0 4px 15px rgba(0,0,0,0.3);
  }
  .title {
    color: #00e676;
    font-size: 16px;
    margin-bottom: 12px;
    font-weight: bold;
  }
  .zikr {
    font-size: 20px;
    color: #ffffff;
    min-height: 60px;
    display: flex;
    align-items: center;
    justify-content: center;
    line-height: 1.5;
    margin-bottom: 10px;
  }
  .sub {
    font-size: 12px;
    color: #8b9bb4;
  }
</style>
</head>
<body>
  <div class="card">
    <div class="title">✨ جاري معالجة المستند واستخراج النصوص... استغل الوقت بالذكر</div>
    <div class="zikr" id="zikr-box">سُبْحَانَ اللَّهِ وَبِحَمْدِهِ ، سُبْحَانَ اللَّهِ الْعَظِيمِ</div>
    <div class="sub">يتغيّر الذكر تلقائياً كل بضع ثوانٍ</div>
  </div>

<script>
  const azkar = [
    "سُبْحَانَ اللَّهِ وَبِحَمْدِهِ ، سُبْحَانَ اللَّهِ الْعَظِيمِ",
    "أَسْتَغْفِرُ اللَّهَ الْعَظِيمَ وَأَتُوبُ إِلَيْهِ",
    "لا حَوْلَ وَلا قُوَّةَ إِلاَّ بِاللَّهِ الْعَلِيِّ الْعَظِيمِ",
    "اللَّهُمَّ صَلِّ وَسَلِّمْ وَبَارِكْ عَلَى نَبِيِّنَا مُحَمَّدٍ",
    "لا إِلَهَ إِلاَّ أَنْتَ سُبْحَانَكَ إِنِّي كُنْتُ مِنَ الظَّالِمِينَ",
    "الْحَمْدُ لِلَّهِ رَبِّ الْعَالَمِينَ",
    "سُبْحَانَ اللَّهِ ، وَالْحَمْدُ لِلَّهِ ، وَلا إِلَهَ إِلاَّ اللَّهُ ، وَاللَّهُ أَكْبَرُ"
  ];
  let idx = 0;
  setInterval(() => {
    idx = (idx + 1) % azkar.length;
    document.getElementById("zikr-box").innerText = azkar[idx];
  }, 4000);
</script>
</body>
</html>
"""


# ----------------------------------------------------
# 2. تنقية الصور واستخراج النصوص المتقدم
# ----------------------------------------------------
def clean_and_enhance_image(img_pil: Image.Image):
    """تحسين تباين الصورة وتنظيف العلامات المائية"""
    if cv2 is not None:
        open_cv_image = np.array(img_pil.convert('RGB')) 
        gray = cv2.cvtColor(open_cv_image, cv2.COLOR_RGB2GRAY)
        clahe = cv2.createCLAHE(clipLimit=2.5, tileGridSize=(8,8))
        return clahe.apply(gray)
    else:
        gray = ImageOps.grayscale(img_pil)
        enhancer = ImageEnhance.Contrast(gray)
        return enhancer.enhance(2.0)


def process_deep_ocr(img_pil: Image.Image, translator_func) -> str:
    """محاولة الاستخراج متعددة المستويات (EasyOCR -> Tesseract)"""
    extracted_text = ""
    
    # 1. التجربة الأولى: عبر EasyOCR
    reader = load_easyocr_reader()
    if reader is not None:
        try:
            processed_img = clean_and_enhance_image(img_pil)
            results = reader.readtext(processed_img, detail=0, paragraph=True)
            extracted_text = "\n".join(results).strip()
        except Exception:
            extracted_text = ""

    # 2. التجربة الثانية: التراجع إلى Tesseract إذا فشل EasyOCR
    if not extracted_text and pytesseract is not None:
        try:
            processed_img = clean_and_enhance_image(img_pil)
            extracted_text = pytesseract.image_to_string(processed_img, lang="fra+eng", config=r'--oem 3 --psm 6').strip()
        except Exception:
            extracted_text = ""

    if extracted_text and len(extracted_text) >= 3:
        return translator_func(extracted_text)
    
    return "[الصفحة تحتوي على مخططات/رسوم توضيحية بدون نص قابل للقراءة]"


# ----------------------------------------------------
# 3. استخراج الصفحات وبناء Side-by-Side PDF
# ----------------------------------------------------
def extract_pages_from_file(uploaded_file):
    file_ext = uploaded_file.name.split(".")[-1].lower()
    pages_data = []

    if file_ext == "pdf":
        uploaded_file.seek(0)
        doc = fitz.open(stream=uploaded_file.read(), filetype="pdf")
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text("text").strip()
            
            # التقاط صورة عالية الوضوح 300 DPI
            pix = page.get_pixmap(dpi=300)
            img_pil = Image.open(io.BytesIO(pix.tobytes("png")))
            w, h = page.rect.width, page.rect.height
            pages_data.append((page_num, text, img_pil, w, h))
        doc.close()

    elif file_ext in ["pptx", "ppt"]:
        uploaded_file.seek(0)
        prs = Presentation(uploaded_file)
        w, h = 720, 540
        for idx, slide in enumerate(prs.slides):
            text_runs = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            text_runs.append(paragraph.text.strip())
            full_text = "\n".join(text_runs).strip()
            
            # إنشاء خلفية للصورة في PPTX
            img_pil = Image.new('RGB', (int(w), int(h)), color=(245, 247, 250))
            pages_data.append((idx, full_text, img_pil, w, h))

    return pages_data


def generate_side_by_side_pdf_safe(uploaded_file, translator_func):
    pages_raw = extract_pages_from_file(uploaded_file)
    total_pages = len(pages_raw)

    progress_bar = st.progress(0)
    status_text = st.empty()
    status_text.text("🚀 جاري استخراج النصوص وترجمتها...")

    results = []
    for completed, (p_num, text, img_pil, w, h) in enumerate(pages_raw, start=1):
        # 1. الاستخراج المباشر للنص الرقمي
        if text and len(text) >= 10:
            translated_text = translator_func(text)
        else:
            # 2. الاستخراج العميق عبر OCR للصورة
            translated_text = process_deep_ocr(img_pil, translator_func)

        results.append((p_num, translated_text, img_pil, w, h))
        progress_bar.progress(completed / total_pages)
        status_text.text(f"⚡ تم ترجمة {completed} من أصل {total_pages} صفحات...")

    status_text.text("🎨 جاري إنشاء الـ PDF النهائي المقسوم وتطبيق العلامة المائية...")

    new_doc = fitz.open()
    all_translated_texts = []

    for page_num, translated_text, final_img_pil, orig_w, orig_h in results:
        all_translated_texts.append(f"--- Page/Slide {page_num + 1} ---\n{translated_text}")

        buffer = io.BytesIO()
        doc_temp = SimpleDocTemplate(
            buffer,
            pagesize=(orig_w, orig_h),
            rightMargin=15, leftMargin=15, topMargin=15, bottomMargin=15,
        )

        styles = getSampleStyleSheet()
        watermark_style = ParagraphStyle(
            "WatermarkStyle", parent=styles["Normal"],
            fontName="Helvetica-Bold", fontSize=8, textColor="#ff4b4b", alignment=1, spaceAfter=2
        )
        title_style = ParagraphStyle(
            "SideTitleStyle", parent=styles["Heading2"],
            fontName="Helvetica-Bold", fontSize=9, spaceAfter=4
        )

        char_count = len(translated_text)
        f_size = 6 if char_count > 1500 else (7 if char_count > 800 else 8)
        leading = f_size + 2

        dynamic_style = ParagraphStyle(
            "DynamicStyle", parent=styles["Normal"],
            fontName="Helvetica", fontSize=f_size, leading=leading, spaceAfter=3
        )

        story = [
            Paragraph("— TRANSLATED BY MEDICAL TRANSLATOR AGENT —", watermark_style),
            Paragraph("BY CHARFI DEKRA", watermark_style),
            Spacer(1, 3),
            Paragraph(f"--- Translation Page/Slide {page_num + 1} ---", title_style)
        ]

        for para in translated_text.split("\n\n"):
            if para.strip():
                formatted = para.strip().replace("\n", "<br/>")
                story.append(Paragraph(formatted, dynamic_style))
                story.append(Spacer(1, 2))

        doc_temp.build(story)
        buffer.seek(0)

        total_width = orig_w * 2
        combo_page = new_doc.new_page(width=total_width, height=orig_h)

        # الصفحة الأصلية (يسار)
        img_byte_arr = io.BytesIO()
        final_img_pil.save(img_byte_arr, format="PNG")
        combo_page.insert_image(
            fitz.Rect(0, 0, orig_w, orig_h),
            stream=img_byte_arr.getvalue()
        )

        # النص المترجم (يمين)
        translated_pdf_doc = fitz.open(stream=buffer.getvalue(), filetype="pdf")
        combo_page.show_pdf_page(
            fitz.Rect(orig_w, 0, total_width, orig_h),
            translated_pdf_doc, 0
        )
        gc.collect()

    output_buffer = io.BytesIO()
    new_doc.save(output_buffer)
    new_doc.close()

    progress_bar.empty()
    status_text.empty()

    output_buffer.seek(0)
    full_text_combined = "\n\n".join(all_translated_texts)
    return output_buffer.getvalue(), full_text_combined


# ----------------------------------------------------
# 4. واجهة Streamlit الرئيسية
# ----------------------------------------------------
st.set_page_config(page_title="MEDICAL TRANSLATOR AGENT", page_icon="🩺", layout="wide")

st.title("🩺 MEDICAL TRANSLATOR AGENT")
st.caption("Advanced Medical Translation Engine | **BY CHARFI DEKRA**")

tab_text, tab_file = st.tabs(["📝 ترجمة نص مباشر", "📄 ترجمة ملف (PDF, Scanned PDF, PowerPoint)"])

with tab_text:
    st.subheader("ترجمة النص الطبي المباشر وتصحيح المصطلحات")
    user_input_text = st.text_area(
        label="أدخل النص المراد ترجمته (فرنسي / عربي):",
        height=200,
        placeholder="أكتب أو ألصق النص هنا...",
    )

    if st.button("ترجمة النص", key="btn_translate_text"):
        if user_input_text.strip():
            with st.spinner("جاري الترجمة والتصحيح الأكاديمي..."):
                result = translate_document(user_input_text)
                st.success("✅ تمت الترجمة بنجاح!")
                st.text_area(label="النص المترجم والمصحح:", value=result, height=250)
        else:
            st.warning("يرجى إدخال نص أولاً.")

with tab_file:
    st.subheader("رفع وترجمة الملفات (يدعم المستندات الممسوحة ضوئياً)")
    uploaded_file = st.file_uploader(
        "قم برفع الملف (PDF عادي، PDF ممسوح ضوئياً، PowerPoint .pptx)", 
        type=["pdf", "pptx", "ppt"]
    )

    if uploaded_file is not None:
        st.success("تم استلام الملف بنجاح!")

        if st.button("ترجمة المستند وتوليد PDF المقسوم", key="btn_translate_file"):
            azkar_container = st.empty()
            with azkar_container.container():
                components.html(AZKAR_HTML, height=200)

            try:
                final_pdf_bytes, combined_text = generate_side_by_side_pdf_safe(
                    uploaded_file, translate_document
                )
                
                azkar_container.empty()
                st.success("🎉 اكتملت المعالجة والترجمة بنجاح!")

                st.text_area(label="معاينة النص المترجم والمصحح:", value=combined_text, height=250)

                st.download_button(
                    label="📥 تحميل الملف المترجم المقسوم (PDF) - BY CHARFI DEKRA",
                    data=final_pdf_bytes,
                    file_name=f"translated_BY_CHARFI_DEKRA_{uploaded_file.name.split('.')[0]}.pdf",
                    mime="application/pdf",
                )
            except Exception as e:
                azkar_container.empty()
                st.error(f"حدث خطأ أثناء المعالجة: {e}")
